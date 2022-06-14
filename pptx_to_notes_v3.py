from pptx import Presentation
import os
import openai
from dotenv import load_dotenv
import json
from transformers import GPT2Model, GPT2Config, GPT2Tokenizer

# ? Enter your file name here, as contained in the inputs folder

file_name = "test_1"

# ! PowerPoint SETUP

# Powerpoint to read from
pptx_path = f"./inputs/{file_name}.pptx"
# file to write to

# ! GPT-3 SETUP

# AI setup and settings
load_dotenv('.env')
api_key = os.environ.get('OPENAI-API-KEY')
openai.api_key = api_key

# config tokenizer
configuration = GPT2Config()
model = GPT2Model(configuration)
tokenizer = GPT2Tokenizer.from_pretrained('gpt2')


def PowerPointtoText():

    if os.path.exists(pptx_path):

        prs = Presentation(pptx_path)

        # pptx setup
        x = 1
        pptx_text = []

        # convert the pptx to text
        for slide in prs.slides:

            pptx_text.append(f"Slide {x}")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_text.append(run.text)
            
            x = x + 1

        return pptx_text

def getResponse(modifier, prompt, temperature):

    prompt = prompt

    # get a response from the AI
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=f'{prompt} \n\n {modifier}',
        temperature=temperature,
        max_tokens=2001,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )

    return f"{response.choices[0].text}"

def IterateOverPrompt(modifier, prompt_full, temperature):

    all_responses = ''
    block_length = 3900
    x = 0
    y = block_length
    count = 1
    prompt_full_length = sum(len(i) for i in prompt_full)

    while prompt_full_length > y:

        print("in the overage loop")

        prompt = prompt_full[x:y]
        response = getResponse(modifier, prompt, temperature)

        all_responses += response
        
        prompt = prompt_full[y:y + block_length]

        y =+ block_length

        print(f"Block {count} indexed")
        count = count + 1

    if prompt_full_length < block_length:

        print("in the regular loop")

        prompt = prompt_full
        all_responses += getResponse(modifier, prompt, temperature)

        print(f"Block {count} indexed")

    return all_responses

def WriteFile(file_name, to_write):
    # write the file
    with open(file_name, 'w') as f:
        json_final = json.dumps(to_write, indent=4)
        f.write(json_final)

# ! Lets do it

pptx_text = PowerPointtoText()

print("PowerPoint converted to text \n")

print(pptx_text)

key_points = IterateOverPrompt("Key points from each slide:", pptx_text, .66)

print("Key points identified \n")
print(key_points)

q_and_a = IterateOverPrompt(
    "| Question | Answer |" , key_points, .66)

print("Questions and answers generated \n")
print(q_and_a)

WriteFile(f"./outputs/{file_name}.json", q_and_a)