from pptx import Presentation
import os
import openai
from dotenv import load_dotenv
from json import dumps, loads

from transformers import GPT2Model, GPT2Config, GPT2Tokenizer

# Powerpoint to read from
pptx = "/inputs/powerpoint.pptx"
# file to write to
file_name = "pptx_test_1.txt"

# ! SETUP

# AI setup and settings
load_dotenv('.env')
api_key = os.environ.get('OPENAI-API-KEY')
openai.api_key = api_key

# config tokenizer
configuration = GPT2Config()
model = GPT2Model(configuration)
tokenizer = GPT2Tokenizer.from_pretrained('gpt2')

# pptx setup
prs = Presentation(pptx)
text_total = []
total_summary = ""
x = 1
pptx_len = len(prs.slides)
text_runs = []

# convert the pptx to text

for slide in prs.slides:

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

output_file_name = f"{file_name.removesuffix('.txt')}_summary"

def read_file(file_path):
    with open(file_path, 'r') as f:
        return f.read()

def get_response():
    # get a response from the AI
    return openai.Completion.create(
    engine="text-davinci-002",
    prompt = f'{modifier} {prompt}',
    temperature = temperature,
    max_tokens = 1000,
    top_p=1,
    frequency_penalty=1,
    presence_penalty=1
    )

all_responses = ''
block_length = 3900
x = 0
y = block_length
count = 1

# modifiers
modifier = "write this text to a json file of key terms : definitions"

# ? Change temp here
# temp
temperature = 0

# set prompt to our converted pptx
prompt_full = text_runs

while y <= len(prompt_full):
    
    prompt = prompt_full[x:y]
    tokens = len(tokenizer(prompt)["attention_mask"])

    while tokens > 1024:
        overage = tokens - 1024
        y = y - (overage * 4)
        prompt = prompt_full[x:y]
        tokens = len(tokenizer(prompt)["attention_mask"])

    # get the response
    response = get_response()

    # add the response to the output
    all_responses += f"{response.choices[0].text}"

    # move to the next block
    x = y
    y = y + block_length

    # let the user know what block we're on
    print(f"Block {count} indexed")
    count = count + 1

# ! AUTO FILE NAMING
# counter for file naming (stored in JSON file)
if os.path.exists("counter.json"):
    counter_int = int(loads(open("counter.json", "r").read()))
else:
    counter_int=1
print(f"Counter: {counter_int}")

# ! WRITE FILE
with open(f"article-summaries/{output_file_name}_{counter_int}.txt", 'w') as f:
    # call write function
    f.write(all_responses)

# if we write a file, update the counter for naming of the next file
counter_int += 1
with open("counter.json", "w") as f:
    f.write(dumps(str(counter_int)))