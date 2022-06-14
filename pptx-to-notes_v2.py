from pptx import Presentation
import os
import openai
from dotenv import load_dotenv
import json
from transformers import GPT2Model, GPT2Config, GPT2Tokenizer, PerceiverForImageClassificationConvProcessing

from rework_iteration_function import dosomething

# ? Enter your file name here, as contained in the inputs folder

file_name = "test_1"

# ! PowerPoint SETUP

# Powerpoint to read from
pptx_path = f"./inputs/{file_name}.pptx"
# file to write to

# ! GPT-3 SETUP

# AI setup and settings
load_dotenv('.env')
api_key = os.environ.get('OPENAI-API-KEY')
openai.api_key = api_key

# config tokenizer
configuration = GPT2Config()
model = GPT2Model(configuration)
tokenizer = GPT2Tokenizer.from_pretrained('gpt2')


def PowerPointtoText():

    if os.path.exists(pptx_path):

        prs = Presentation(pptx_path)

        # pptx setup
        x = 1
        pptx_text = []

        # convert the pptx to text
        for slide in prs.slides:

            pptx_text.append(f"Slide {x}")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_text.append(run.text)
            
            x = x + 1

        return pptx_text


def get_response(modifier, prompt, temperature):

    prompt = prompt

    # get a response from the AI
    return openai.Completion.create(
        engine="text-davinci-002",
        prompt=f'{modifier} {prompt}',
        temperature=temperature,
        max_tokens=1000,
        top_p=1,
        frequency_penalty=1,
        presence_penalty=1
    )

                        # limit = 3
                        # def dosomething():
                        #     return "something"
                        # x = 0
                        # y = limit

                        # while string > limit:
                        #     slice = string[x:y]
                        #     dosomething(slice)
                        #     string = string[y:]

                        # if len(string) < limit:
                        #     dosomething(string)

def IterateOverPrompt(modifier, prompt_full, temperature):

    all_responses = ''
    block_length = 3900
    x = 0
    y = block_length
    count = 1
    prompt_full_length = sum(len(i) for i in prompt_full)

    if prompt_full_length < y:
        go = True

    print(f"length of pptx text is {prompt_full_length} ")
    print(f"y is {y}")

    # I had to use this term below bexause ["ben", "sam", "Jimmy"] computes to 3 for length
    # this term is the aggregate length of the list of strings


    while y <= prompt_full_length or y>0:

        print("in the loop")

        prompt = prompt_full[x:y]
        tokens = len(tokenizer(prompt)["attention_mask"])

        while tokens > 1024:
            overage = tokens - 1024
            y = y - (overage * 4)
            prompt = prompt_full[x:y]
            tokens = len(tokenizer(prompt)["attention_mask"])

        print(f"prompt is from {x} to {y}")

        # get the response
        response = get_response(modifier, prompt, temperature)

        # add the response to the output
        all_responses += f"{response.choices[0].text}"

        # move to the next block
        x = y
        y = y + block_length

        print(f"y is now {y}")

        # let the user know what block we're on
        print(f"Block {count} indexed")
        count = count + 1


    return all_responses


def WriteFile(file_name, to_write):
    # write the file
    with open(file_name, 'w') as f:
        json_final = json.dumps(to_write)
        f.write(json_final)

# ! Lets do it

pptx_text = PowerPointtoText()

print("PowerPoint converted to text")
print(pptx_text)

key_points = IterateOverPrompt("Identify key points from each slide, in JSON format", pptx_text, .7)

print("Key points identified")
print(key_points)

q_and_a = IterateOverPrompt(
    "Identify questions and answers from each slide, in JSON format", pptx_text, .7)

print("Questions and answers generated")
print(q_and_a)

exit()
WriteFile(f"./outputs/{file_name}.json", q_and_a)